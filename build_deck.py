"""Build the Cerebro overview deck — Howden branded.

Brand tokens sourced from colors_and_type.css:
  moss-green  #173F35  (primary dark, brand)
  pistachio   #E0EC89  (primary accent, Cerebro logotype)
  nero        #1A1A1A  (primary text)
  dark-gray   #666666  (secondary text)
  soft-white  #F6F6F6  (panel background)
  cobalt      #0857C3  (action)
  mustard     #FFBF3F  (CTA)

  Display serif  GT Ultra Median
  Sans           Aktiv Grotesk
  Mono           IBM Plex Mono
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Howden brand palette ----
MOSS        = RGBColor(0x17, 0x3F, 0x35)
MOSS_85     = RGBColor(0x2A, 0x52, 0x48)   # slight lift for dividers on dark bg
PISTACHIO   = RGBColor(0xE0, 0xEC, 0x89)
PISTACHIO_L = RGBColor(0xF9, 0xFB, 0xE7)
NERO        = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY   = RGBColor(0x66, 0x66, 0x66)
SOFT_WHITE  = RGBColor(0xF6, 0xF6, 0xF6)
CLOUD_GRAY  = RGBColor(0xE5, 0xE5, 0xE5)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
COBALT      = RGBColor(0x08, 0x57, 0xC3)
MUSTARD     = RGBColor(0xFF, 0xBF, 0x3F)

# ---- Typography ----
F_SERIF = "GT Ultra Median"
F_SANS  = "Aktiv Grotesk"
F_MONO  = "IBM Plex Mono"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

def bg(slide, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.line.fill.background()
    r.fill.solid(); r.fill.fore_color.rgb = color
    slide.shapes._spTree.remove(r._element); slide.shapes._spTree.insert(2, r._element)
    return r

def txt(slide, x, y, w, h, text, *, size=18, bold=False, color=NERO, font=F_SANS,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, tracking=0, line_space=1.15, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_space
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        if tracking:
            r._r.get_or_add_rPr().set("spc", str(tracking))
    return box

def eyebrow(slide, x, y, text, color=MOSS):
    return txt(slide, x, y, Inches(8), Inches(0.3), text.upper(),
               size=10, bold=True, color=color, font=F_MONO, tracking=400)

def card(slide, x, y, w, h, *, fill=WHITE, border=CLOUD_GRAY, radius=0.06, border_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = radius
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border
        shp.line.width = Pt(border_w)
    shp.shadow.inherit = False
    return shp

def hr(slide, x, y, w, color=CLOUD_GRAY, weight=0.75):
    ln = slide.shapes.add_connector(1, x, y, x + w, y)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln

def brand_lockup(slide, x, y, *, dark_bg=False):
    """Render the Cerebro + 'from HOWDEN' lockup exactly like the app header."""
    # "Cerebro" in pistachio serif
    txt(slide, x, y, Inches(3), Inches(0.55),
        "Cerebro", size=22, bold=False, color=PISTACHIO if dark_bg else MOSS,
        font=F_SERIF, line_space=1.0)
    # "from HOWDEN" tag, small sans, next to it
    fy = y + Inches(0.18)
    txt(slide, x + Inches(1.4), fy, Inches(0.6), Inches(0.3),
        "from", size=8, bold=False,
        color=WHITE if dark_bg else DARK_GRAY, font=F_SANS, tracking=100)
    txt(slide, x + Inches(1.7), fy, Inches(1.4), Inches(0.3),
        "HOWDEN", size=9, bold=True,
        color=WHITE if dark_bg else NERO, font=F_SANS, tracking=100)

def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# Slide 1 · Cover — moss green with pistachio Cerebro logotype
# ============================================================
s = new_slide(); bg(s, MOSS)
# Brand lockup top-left
brand_lockup(s, Inches(0.75), Inches(0.6), dark_bg=True)

# Small eyebrow
txt(s, Inches(0.75), Inches(2.5), Inches(6), Inches(0.3),
    "AN INTERNAL PROTOTYPE", size=10, bold=True, color=PISTACHIO,
    font=F_MONO, tracking=400)

# Hero line — serif, large
txt(s, Inches(0.75), Inches(3.0), Inches(12), Inches(1.8),
    "Triage at the speed", size=72, color=WHITE, font=F_SERIF, line_space=1.0)
txt(s, Inches(0.75), Inches(4.05), Inches(12), Inches(1.8),
    "of email.", size=72, color=PISTACHIO, font=F_SERIF, line_space=1.0)

# Footer — small mono
txt(s, Inches(0.75), Inches(6.7), Inches(12), Inches(0.4),
    "2026 · OVERVIEW", size=9, bold=True, color=PISTACHIO,
    font=F_MONO, tracking=400)

# ============================================================
# Slide 2 · The problem
# ============================================================
s = new_slide(); bg(s, SOFT_WHITE)
brand_lockup(s, Inches(0.75), Inches(0.55))
eyebrow(s, Inches(0.75), Inches(1.5), "Today")

txt(s, Inches(0.75), Inches(2.0), Inches(12), Inches(1.6),
    "Every referral needs a human", size=52, color=NERO, font=F_SERIF, line_space=1.05)
txt(s, Inches(0.75), Inches(2.95), Inches(12), Inches(1.6),
    "to read, classify, and route.", size=52, color=DARK_GRAY, font=F_SERIF, line_space=1.05)

hr(s, Inches(0.75), Inches(5.6), Inches(11.8))
txt(s, Inches(0.75), Inches(5.85), Inches(12), Inches(1.2),
    "~15 minutes each. Hundreds a day. The queue never catches up.",
    size=18, color=NERO, font=F_SANS, line_space=1.35)

# ============================================================
# Slide 3 · What Cerebro does
# ============================================================
s = new_slide(); bg(s, SOFT_WHITE)
brand_lockup(s, Inches(0.75), Inches(0.55))
eyebrow(s, Inches(0.75), Inches(1.5), "With Cerebro")

txt(s, Inches(0.75), Inches(2.0), Inches(12), Inches(1.4),
    "It reads the email.", size=52, color=NERO, font=F_SERIF, line_space=1.05)
txt(s, Inches(0.75), Inches(2.95), Inches(12), Inches(1.4),
    "Routes the risk.", size=52, color=NERO, font=F_SERIF, line_space=1.05)
txt(s, Inches(0.75), Inches(3.9), Inches(12), Inches(1.4),
    "In under a second.", size=52, color=MOSS, font=F_SERIF, line_space=1.05)

hr(s, Inches(0.75), Inches(5.9), Inches(11.8))
txt(s, Inches(0.75), Inches(6.15), Inches(12), Inches(1.0),
    "Claude extracts the risk. Rules pick the platform. Brokers see the why.",
    size=16, color=DARK_GRAY, font=F_SANS, line_space=1.35)

# ============================================================
# Slide 4 · How it works (3 bento)
# ============================================================
s = new_slide(); bg(s, SOFT_WHITE)
brand_lockup(s, Inches(0.75), Inches(0.55))
eyebrow(s, Inches(0.75), Inches(1.5), "How it works")
txt(s, Inches(0.75), Inches(1.85), Inches(12), Inches(0.9),
    "Three steps. Always traceable.", size=34, color=NERO, font=F_SERIF, line_space=1.05)

bento_y = Inches(3.15)
bento_h = Inches(3.6)
gap = Inches(0.22)
total_w = SW - Inches(1.5)
box_w = (total_w - gap*2) / 3

steps = [
    ("01", "INGEST",  "Email, slip, SOV, or pasted thread.",
     "Claude normalises any source into a structured submission."),
    ("02", "TRIAGE",  "Extract class, premium, geography, losses.",
     "Every field carries a confidence score and source span."),
    ("03", "ROUTE",   "Nine rules. Six destinations. One decision.",
     "xTrade · WhiteSpace · HAT · GXB · Acturis · Manual review."),
]
for i, (num, title, lead, body) in enumerate(steps):
    x = Inches(0.75) + (box_w + gap) * i
    card(s, x, bento_y, box_w, bento_h, fill=WHITE, border=CLOUD_GRAY)
    txt(s, x + Inches(0.35), bento_y + Inches(0.35), box_w, Inches(0.35),
        num, size=10, bold=True, color=MOSS, font=F_MONO, tracking=400)
    txt(s, x + Inches(0.35), bento_y + Inches(0.75), box_w, Inches(0.5),
        title, size=13, bold=True, color=NERO, font=F_SANS, tracking=200)
    txt(s, x + Inches(0.35), bento_y + Inches(1.5), box_w - Inches(0.7), Inches(1.2),
        lead, size=22, color=NERO, font=F_SERIF, line_space=1.2)
    txt(s, x + Inches(0.35), bento_y + Inches(2.7), box_w - Inches(0.7), Inches(1.0),
        body, size=12, color=DARK_GRAY, font=F_SANS, line_space=1.4)

# ============================================================
# Slide 5 · What's inside (2x2 bento)
# ============================================================
s = new_slide(); bg(s, SOFT_WHITE)
brand_lockup(s, Inches(0.75), Inches(0.55))
eyebrow(s, Inches(0.75), Inches(1.5), "Inside")
txt(s, Inches(0.75), Inches(1.85), Inches(12), Inches(0.9),
    "Four things working together.", size=34, color=NERO, font=F_SERIF, line_space=1.05)

grid_y = Inches(3.1)
grid_w = SW - Inches(1.5)
grid_h = Inches(3.9)
gap = Inches(0.22)
cell_w = (grid_w - gap) / 2
cell_h = (grid_h - gap) / 2

cells = [
    ("Extraction",   "Claude reads the email.",
     "Structured fields with per-field confidence. Unsure cases flagged for review.",
     False),
    ("Rules",        "Nine routing rules.",
     "Sanctions gates, binder matches, class, size, loss history. Priority-ordered.",
     False),
    ("Destinations", "Six placement platforms.",
     "xTrade · WhiteSpace PPL · HAT · GXB · Acturis · Manual review.",
     True),  # feature card — moss-green with pistachio
    ("Audit",        "Every decision, traceable.",
     "Which rule fired, which didn't, why. Broker overrides feed back into learning.",
     False),
]
for i, (eb, lead, body, feature) in enumerate(cells):
    col, row = i % 2, i // 2
    x = Inches(0.75) + (cell_w + gap) * col
    y = grid_y + (cell_h + gap) * row
    if feature:
        card(s, x, y, cell_w, cell_h, fill=MOSS, border=None)
        eb_col = PISTACHIO; lead_col = WHITE; body_col = PISTACHIO_L
    else:
        card(s, x, y, cell_w, cell_h, fill=WHITE, border=CLOUD_GRAY)
        eb_col = MOSS; lead_col = NERO; body_col = DARK_GRAY
    txt(s, x + Inches(0.4), y + Inches(0.35), cell_w, Inches(0.35),
        eb.upper(), size=10, bold=True, color=eb_col, font=F_MONO, tracking=400)
    txt(s, x + Inches(0.4), y + Inches(0.8), cell_w - Inches(0.8), Inches(0.8),
        lead, size=26, color=lead_col, font=F_SERIF, line_space=1.15)
    txt(s, x + Inches(0.4), y + cell_h - Inches(1.0), cell_w - Inches(0.8), Inches(0.9),
        body, size=13, color=body_col, font=F_SANS, line_space=1.45)

# ============================================================
# Slide 6 · Where we are / preview
# ============================================================
s = new_slide(); bg(s, SOFT_WHITE)
brand_lockup(s, Inches(0.75), Inches(0.55))
eyebrow(s, Inches(0.75), Inches(1.5), "Where we are")

txt(s, Inches(0.75), Inches(2.0), Inches(12), Inches(1.4),
    "Live prototype.", size=60, color=NERO, font=F_SERIF, line_space=1.0)
txt(s, Inches(0.75), Inches(3.1), Inches(12), Inches(1.4),
    "Preview today.", size=60, color=MOSS, font=F_SERIF, line_space=1.0)

# Preview card — pistachio tint, moss border for brand cohesion
card_y = Inches(5.0)
card(s, Inches(0.75), card_y, SW - Inches(1.5), Inches(1.6),
     fill=PISTACHIO_L, border=PISTACHIO, border_w=1.25)
txt(s, Inches(1.1), card_y + Inches(0.3), Inches(10), Inches(0.35),
    "PREVIEW", size=10, bold=True, color=MOSS, font=F_MONO, tracking=400)
txt(s, Inches(1.1), card_y + Inches(0.65), Inches(11), Inches(0.7),
    "lukataylo.github.io/Cerebro",
    size=26, bold=True, color=NERO, font=F_MONO, tracking=-50)

# Footer
txt(s, Inches(0.75), Inches(7.0), Inches(12), Inches(0.4),
    "Admin console · Broker workbench · Routing flow · Rules · Audit",
    size=9, bold=True, color=DARK_GRAY, font=F_MONO, tracking=300)

out = "/Users/lukadadiani/Documents/Cerebro/Cerebro.pptx"
prs.save(out)
print(f"wrote {out}")
